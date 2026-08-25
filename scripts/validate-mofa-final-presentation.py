#!/usr/bin/env python3
from __future__ import annotations

import hashlib
import json
import re
import subprocess
from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "artifacts/submission/2026-mofa-ai/final-presentation"
PPTX = OUT / "claimgate-oda-final-presentation.pptx"
PDF = OUT / "claimgate-oda-final-presentation.pdf"
NOTES = OUT / "claimgate-oda-speaker-notes.md"
SOURCE = OUT / "claimgate-oda-final-presentation-source.md"
MANIFEST = OUT / "SUBMISSION-MANIFEST.json"
SUMS = OUT / "SHA256SUMS.txt"
PACKAGE = OUT / "claimgate-oda-final-presentation-package.zip"


def digest(path: Path) -> str:
    value = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            value.update(chunk)
    return value.hexdigest()


def require(condition: bool, message: str) -> None:
    if not condition:
        raise SystemExit(f"final-presentation: FAIL: {message}")


def main() -> None:
    required = [PPTX, PDF, NOTES, SOURCE, MANIFEST, SUMS, PACKAGE]
    for path in required:
        require(path.is_file() and path.stat().st_size > 0, f"missing or empty {path.relative_to(ROOT)}")

    deck = Presentation(PPTX)
    require(len(deck.slides) == 10, f"slide count={len(deck.slides)}, want 10")
    require(abs(deck.slide_width / deck.slide_height - 16 / 9) < 0.02, "deck is not 16:9")

    pdf_info = subprocess.check_output(["pdfinfo", str(PDF)], text=True)
    require(re.search(r"^Pages:\s+10$", pdf_info, re.MULTILINE) is not None, "PDF must have 10 pages")
    pdf_text = subprocess.check_output(["pdftotext", str(PDF), "-"], text=True)

    combined = "\n".join([SOURCE.read_text(), NOTES.read_text(), pdf_text])
    for phrase in [
        "No Anchor, No Claim",
        "AI Curator, Not Judge",
        "Evidence Pack",
        "외교부_국가별 안전정보",
        "offline fixture",
    ]:
        require(phrase.lower() in combined.lower(), f"required phrase missing: {phrase}")

    for pattern in [
        r"010-\d{3,4}-\d{4}",
        r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}",
        r"hallucination을 제거",
        r"자동으로 진실을 판정",
        r"live OpenAPI 연동 완료",
        r"production accuracy 검증 완료",
    ]:
        require(re.search(pattern, combined, re.IGNORECASE) is None, f"forbidden/private claim matched: {pattern}")

    for name in ["01-guided-demo-start.png", "02-review-workspace.png"]:
        image = Image.open(OUT / "assets" / name)
        require(image.width >= 1400 and image.height >= 900, f"presentation image too small: {name}")

    manifest = json.loads(MANIFEST.read_text())
    require(manifest.get("schema") == "claimgate.mofa-final-presentation/v1", "unexpected manifest schema")
    require(manifest.get("slideCount") == 10, "manifest slideCount must be 10")
    for item in manifest.get("files", []):
        path = OUT / item["path"]
        require(path.is_file(), f"manifest path missing: {item['path']}")
        require(path.stat().st_size == item["bytes"], f"size mismatch: {item['path']}")
        require(digest(path) == item["sha256"], f"checksum mismatch: {item['path']}")

    expected_zip = {
        "claimgate-oda-final-presentation.pdf",
        "claimgate-oda-final-presentation.pptx",
        "claimgate-oda-speaker-notes.md",
        "README.md",
        "SHA256SUMS.txt",
        "SUBMISSION-MANIFEST.json",
    }
    with ZipFile(PACKAGE) as archive:
        require(set(archive.namelist()) == expected_zip, "unexpected ZIP member set")
        require(archive.testzip() is None, "ZIP CRC failure")

    print("final-presentation: PASS — 10 slides, PDF/PPTX, notes, checksums, ZIP and claim boundaries verified")


if __name__ == "__main__":
    main()
