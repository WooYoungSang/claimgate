#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "artifacts/submission/2026-mofa-ai/final-presentation"
ASSETS = OUT / "assets"
PPTX_PATH = OUT / "claimgate-oda-final-presentation.pptx"

FONT = "Noto Sans CJK KR"
GREEN = RGBColor(15, 77, 54)
GREEN_2 = RGBColor(20, 115, 79)
MINT = RGBColor(229, 242, 234)
CREAM = RGBColor(247, 249, 246)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(24, 52, 42)
MUTED = RGBColor(91, 111, 101)
LINE = RGBColor(210, 223, 215)
AMBER = RGBColor(229, 158, 53)
RED = RGBColor(205, 82, 74)
PALE_RED = RGBColor(250, 234, 232)
PALE_AMBER = RGBColor(252, 243, 224)


def set_text(run, size: float, color=DARK, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    set_text(run, size, color, bold)
    return box


def add_rich_lines(slide, x, y, w, h, lines, size=17, gap=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.04)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    for index, (text, color, bold) in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = text
        set_text(run, size, color, bold)
    return box


def rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def pill(slide, x, y, w, text, fill=MINT, color=GREEN, size=10.5):
    shape = rect(slide, x, y, w, 0.34, fill=fill, line=fill)
    add_text(slide, x, y + 0.01, w, 0.29, text, size=size, color=color, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    return shape


def title(slide, section, heading, sub=None, page=1, dark=False):
    color = WHITE if dark else DARK
    muted = RGBColor(205, 225, 214) if dark else MUTED
    add_text(slide, 0.62, 0.38, 5.2, 0.28, section.upper(), 10.5, muted, True)
    add_text(slide, 0.62, 0.78, 11.9, 0.62, heading, 27, color, True)
    if sub:
        add_text(slide, 0.64, 1.42, 11.5, 0.38, sub, 12.5, muted)
    add_text(slide, 12.2, 0.40, 0.48, 0.24, f"{page:02d}", 10.5, muted, True,
             align=PP_ALIGN.RIGHT)


def footer(slide, text="CLAIMGATE ODA · 2026 외교 공공데이터·AI 활용 경진대회", dark=False):
    color = RGBColor(188, 211, 198) if dark else RGBColor(127, 148, 137)
    add_text(slide, 0.62, 7.18, 8.2, 0.20, text, 8.4, color)


def add_bullet_card(slide, x, y, w, h, number, heading, body, accent=GREEN):
    rect(slide, x, y, w, h, WHITE, LINE)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(y + 0.22), Inches(0.42), Inches(0.42))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent; circle.line.color.rgb = accent
    add_text(slide, x + 0.22, y + 0.23, 0.42, 0.35, str(number), 11, WHITE, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    add_text(slide, x + 0.78, y + 0.19, w - 1.0, 0.35, heading, 15, DARK, True)
    add_text(slide, x + 0.78, y + 0.61, w - 1.02, h - 0.76, body, 11.3, MUTED)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = GREEN
    add_text(slide, 0.72, 0.66, 4.2, 0.28, "PRODUCT / SERVICE DEVELOPMENT", 10.5, RGBColor(190, 219, 203), True)
    add_text(slide, 0.72, 1.18, 5.2, 0.62, "ClaimGate ODA", 34, WHITE, True)
    add_text(slide, 0.72, 1.92, 5.3, 1.2, "AI가 쓴 문장을\n근거와 책임이 있는 주장으로", 25, WHITE, True)
    add_text(slide, 0.75, 3.37, 4.8, 0.52, "공식 데이터의 근거 위치 · 결정론적 위험 규칙 · 사람의 최종 판단", 12.5, RGBColor(210, 229, 217))
    pill(slide, 0.75, 4.20, 1.82, "NO ANCHOR", RGBColor(36, 97, 72), WHITE, 10)
    pill(slide, 2.68, 4.20, 1.62, "NO CLAIM", AMBER, DARK, 10)
    add_text(slide, 0.75, 5.02, 4.9, 0.52, "2026 외교 공공데이터·AI 활용 경진대회", 12, WHITE, True)
    add_text(slide, 0.75, 5.49, 4.5, 0.40, "제품·서비스 개발 부문", 11.5, RGBColor(202, 225, 212))
    image = ASSETS / "01-guided-demo-start.png"
    slide.shapes.add_picture(str(image), Inches(6.02), Inches(0.62), width=Inches(6.65), height=Inches(5.88))
    add_text(slide, 6.05, 6.68, 6.55, 0.28, "실제 오프라인 고정 예시 데이터 기반 시제품 화면", 9.5, RGBColor(197, 220, 207), align=PP_ALIGN.RIGHT)
    return slide


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CREAM
    title(slide, "01 · Problem", "빠른 초안과 책임 있는 설명 사이의 빈틈", "문장이 자연스럽다는 사실만으로 공식 데이터와의 일치를 보장할 수 없습니다.", 2)
    quote = rect(slide, 0.62, 2.02, 12.05, 1.05, GREEN, GREEN)
    add_text(slide, 0.94, 2.26, 11.4, 0.48, "“공공기관 설명의 작은 오류도 정책 판단과 대외 신뢰의 위험으로 이어집니다.”", 20, WHITE, True, align=PP_ALIGN.CENTER)
    cards = [
        ("01", "최신성", "국가 안전정보가 최신 외교부 데이터와 일치하는가?", RED),
        ("02", "정합성", "사업 대상국·기관·기간·성과가 공개 사업정보와 맞는가?", AMBER),
        ("03", "책임성", "누가 어떤 근거를 확인하고 무엇을 정정했는가?", GREEN_2),
    ]
    for i, (num, head, body, accent) in enumerate(cards):
        x = 0.62 + i * 4.08
        rect(slide, x, 3.46, 3.82, 2.68, WHITE, LINE)
        pill(slide, x + 0.24, 3.70, 0.58, num, accent, WHITE, 10)
        add_text(slide, x + 0.26, 4.22, 3.25, 0.38, head, 18, DARK, True)
        add_text(slide, x + 0.26, 4.78, 3.28, 0.92, body, 13.2, MUTED)
    footer(slide)
    return slide


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE
    title(slide, "02 · Solution", "생성 이후의 책임 있는 검토를 제품화합니다", "AI는 후보를 빠르게 구성하고, 근거·규칙·사람이 최종 상태를 결정합니다.", 3)
    steps = [
        ("AI 후보", "주장·근거 후보\n제안 전용", RGBColor(73, 90, 115)),
        ("Source Anchor", "행·필드·문장\n근거 위치 연결", GREEN_2),
        ("Rule Trace", "값·기관·날짜\n불일치 추적", AMBER),
        ("Human Review", "검증·정정·기각\n최종 결정", RED),
        ("Evidence Pack", "검토된 주장만\n후속 사용", GREEN),
    ]
    for i, (head, body, accent) in enumerate(steps):
        x = 0.50 + i * 2.57
        rect(slide, x, 2.42, 2.08, 2.52, CREAM, LINE)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.72), Inches(2.67), Inches(0.62), Inches(0.62))
        circle.fill.solid(); circle.fill.fore_color.rgb = accent; circle.line.color.rgb = accent
        add_text(slide, x + 0.10, 3.47, 1.88, 0.38, head, 14, DARK, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.16, 4.00, 1.76, 0.62, body, 11.4, MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text(slide, x + 2.12, 3.42, 0.39, 0.38, "→", 18, GREEN, True, align=PP_ALIGN.CENTER)
    rect(slide, 0.62, 5.45, 12.05, 1.05, MINT, MINT)
    add_text(slide, 0.88, 5.68, 11.5, 0.52, "검토자 판정 전에는 보고서·그래프·근거 묶음 투영이 차단됩니다.", 18, GREEN, True, align=PP_ALIGN.CENTER)
    footer(slide)
    return slide


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = DARK
    title(slide, "03 · Product Demo", "세 가지 위험 수준을 한 화면에서 검토", "실제 MOFA ODA DomainPack 오프라인 시제품", 4, dark=True)
    slide.shapes.add_picture(str(ASSETS / "02-review-workspace.png"), Inches(0.62), Inches(1.85), width=Inches(9.12), height=Inches(4.98))
    items = [
        ("RED", "안전정보 충돌", "외교부 국가별 안전정보와 AI 안전 주장 불일치", PALE_RED, RED),
        ("YELLOW", "사업정보 확인", "KOICA 대상국·기간·기관 추가 검토", PALE_AMBER, AMBER),
        ("GREEN", "정의 일치", "ODA 용어 정의 일치 + sampling 후보", MINT, GREEN_2),
    ]
    for i, (tag, head, body, fill, accent) in enumerate(items):
        y = 2.00 + i * 1.50
        rect(slide, 9.98, y, 2.70, 1.24, fill, fill)
        pill(slide, 10.18, y + 0.16, 0.88, tag, accent, WHITE, 8.7)
        add_text(slide, 10.18, y + 0.55, 2.12, 0.27, head, 12.2, DARK, True)
        add_text(slide, 10.18, y + 0.84, 2.18, 0.31, body, 8.9, MUTED)
    footer(slide, dark=True)
    return slide


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CREAM
    title(slide, "04 · Public Data", "공공데이터를 조회 화면이 아니라 주장 검토 기준점으로", "현재 구현과 향후 연동을 명확히 분리합니다.", 5)
    current = [
        ("DATA-001", "외교부_국가별 안전정보", "국가·지역·공지 시점·위험 내용"),
        ("DATA-002", "KOICA 국가별 협력사업", "대상국·사업명·기간·기관"),
        ("DATA-003", "KOICA ODA 용어사전", "공식 용어와 정의 일치"),
    ]
    add_text(slide, 0.72, 2.02, 5.55, 0.35, "CURRENT · OFFLINE FIXTURE", 11, GREEN, True)
    for i, (tag, head, body) in enumerate(current):
        y = 2.50 + i * 1.12
        rect(slide, 0.68, y, 5.75, 0.92, WHITE, LINE)
        pill(slide, 0.90, y + 0.20, 0.94, tag, GREEN, WHITE, 8.5)
        add_text(slide, 2.02, y + 0.15, 3.95, 0.28, head, 13.2, DARK, True)
        add_text(slide, 2.02, y + 0.51, 3.95, 0.24, body, 10.2, MUTED)
    add_text(slide, 6.86, 2.02, 5.5, 0.35, "NEXT · READ-ONLY ADAPTER", 11, AMBER, True)
    future = [
        ("KF", "공공외교 사업 정보", "사업명·기관·기간·성과 설명"),
        ("KAF", "아프리카 공개데이터", "국가·기관·스타트업·교류협력"),
    ]
    for i, (tag, head, body) in enumerate(future):
        y = 2.50 + i * 1.36
        rect(slide, 6.82, y, 5.82, 1.12, WHITE, LINE)
        pill(slide, 7.06, y + 0.25, 0.74, tag, AMBER, WHITE, 9)
        add_text(slide, 8.00, y + 0.18, 4.10, 0.30, head, 14, DARK, True)
        add_text(slide, 8.00, y + 0.60, 4.18, 0.28, body, 10.5, MUTED)
    rect(slide, 6.82, 5.43, 5.82, 1.00, PALE_AMBER, PALE_AMBER)
    add_text(slide, 7.08, 5.64, 5.27, 0.52, "현재 URL은 provenance입니다.\n실시간 OpenAPI 연동·운영 정확도를 주장하지 않습니다.", 12.1, DARK, True, align=PP_ALIGN.CENTER)
    footer(slide)
    return slide


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE
    title(slide, "05 · Differentiation", "자동 진실판정이 아니라 책임 있는 검토 시스템", "다섯 개의 제품 원칙이 AI 권한과 사람의 책임을 분리합니다.", 6)
    principles = [
        ("AI Curator, Not Judge", "AI는 후보 주장과 근거 후보만 제안", RGBColor(72, 91, 116)),
        ("No Anchor, No Claim", "근거 위치 없는 주장은 완료 상태 불가", GREEN),
        ("Deterministic Risk", "명시적 규칙과 rule trace로 위험 표시", AMBER),
        ("Green Sampling", "낮은 위험의 조용한 오류도 표본 확인", GREEN_2),
        ("Evidence Pack First", "검증·정정된 주장만 재사용 산출물로", RED),
    ]
    positions = [(0.66,2.02,3.92),(4.70,2.02,3.92),(8.74,2.02,3.92),(2.68,4.14,3.92),(6.72,4.14,3.92)]
    for (head, body, accent), (x,y,w) in zip(principles, positions):
        rect(slide, x, y, w, 1.74, CREAM, LINE)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(1.74))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.color.rgb = accent
        add_text(slide, x + 0.32, y + 0.28, w - 0.56, 0.36, head, 15.1, DARK, True)
        add_text(slide, x + 0.32, y + 0.84, w - 0.56, 0.50, body, 11.5, MUTED)
    footer(slide)
    return slide


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CREAM
    title(slide, "06 · Architecture", "도메인은 교체하고, 신뢰 불변식은 고정합니다", "AI·UI·DomainPack이 core의 권한 경계를 넘지 못하도록 분리했습니다.", 7)
    layers = [
        ("Controlled UI", "검토 상태를 표시하지만 숨은 판정 권한 없음", RGBColor(225, 237, 230), GREEN),
        ("DomainPack", "ODA·보건·시민데이터 fixture와 판단 규칙", PALE_AMBER, AMBER),
        ("ClaimGate Core", "Claim · Source Anchor · Risk · Review · Evidence Pack", RGBColor(205, 227, 214), GREEN),
    ]
    for i, (head, body, fill, accent) in enumerate(layers):
        y = 1.98 + i * 1.35
        rect(slide, 0.72 + i*0.30, y, 7.18 - i*0.60, 1.02, fill, fill)
        add_text(slide, 1.02 + i*0.30, y + 0.18, 2.15, 0.32, head, 15, accent, True)
        add_text(slide, 3.20 + i*0.15, y + 0.20, 4.05 - i*0.50, 0.38, body, 11.2, DARK)
    rect(slide, 8.42, 1.98, 4.20, 4.06, DARK, DARK)
    add_text(slide, 8.78, 2.30, 3.50, 0.35, "@claimgate/ai-local", 17, WHITE, True)
    add_text(slide, 8.78, 2.83, 3.48, 0.72, "후보 추출 · local model · RAG provenance", 12.5, RGBColor(201, 224, 211))
    add_text(slide, 8.78, 3.80, 3.46, 0.34, "가질 수 없는 권한", 12, AMBER, True)
    add_rich_lines(slide, 8.78, 4.23, 3.42, 1.40, [
        ("× 사실 검증", WHITE, True),
        ("× 위험 최종 결정", WHITE, True),
        ("× accepted anchor 부착", WHITE, True),
        ("× Evidence Pack 투영", WHITE, True),
    ], 11.2, 5)
    footer(slide)
    return slide


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE
    title(slide, "07 · Evidence", "구현 주장은 재현 가능한 테스트와 산출물에 연결됩니다", "2026-08-25 저장소 검증 기준", 8)
    metrics = [
        ("192", "workspace unit tests", GREEN),
        ("103", "core tests", GREEN_2),
        ("20", "local-AI tests", AMBER),
        ("3", "DomainPacks", RED),
    ]
    for i, (value, label, accent) in enumerate(metrics):
        x = 0.66 + i * 3.08
        rect(slide, x, 2.04, 2.82, 1.70, CREAM, LINE)
        add_text(slide, x + 0.15, 2.30, 2.52, 0.58, value, 30, accent, True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.15, 3.06, 2.52, 0.30, label, 10.5, MUTED, True, align=PP_ALIGN.CENTER)
    gates = [
        "오프라인·결정론적 기본 경로",
        "No Anchor / reviewer / projection negative controls",
        "DomainPack conformance",
        "contest overclaim guard",
        "video preflight · kbctl verify",
    ]
    rect(slide, 0.66, 4.18, 12.00, 1.84, MINT, MINT)
    for i, item in enumerate(gates):
        x = 0.94 + (i % 3) * 3.85
        y = 4.48 + (i // 3) * 0.70
        add_text(slide, x, y, 3.45, 0.34, "✓  " + item, 11.2, GREEN, True)
    add_text(slide, 0.78, 6.35, 11.8, 0.38, "정확도·고객 채택·운영 시간 절감 수치는 아직 주장하지 않습니다.", 12.4, RED, True, align=PP_ALIGN.CENTER)
    footer(slide)
    return slide


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CREAM
    title(slide, "08 · Roadmap", "고위험 문장 검토부터 제한적으로 제품화합니다", "문서 생성 전체를 대체하지 않고 검토 비용과 책임 인계 문제부터 해결합니다.", 9)
    steps = [
        (1, "검토 엔진", "Claim · Anchor · Rule · Review · Evidence Pack 안정화"),
        (2, "ODA 시제품", "외교 공공데이터 fixture와 발표 데모 검증"),
        (3, "제한 파일럿", "검토 시간·재확인 시간·정정 재사용률 측정"),
        (4, "운영 연계", "필요한 범위의 읽기 전용 데이터·AI 어댑터"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.66 + i * 3.08
        rect(slide, x, 2.18, 2.82, 3.42, WHITE, LINE)
        pill(slide, x + 0.22, 2.42, 0.62, f"0{num}", GREEN if num <= 2 else AMBER, WHITE, 9)
        add_text(slide, x + 0.22, 3.08, 2.34, 0.40, head, 16, DARK, True)
        add_text(slide, x + 0.22, 3.72, 2.34, 1.06, body, 11.4, MUTED)
        add_text(slide, x + 0.22, 5.04, 2.34, 0.28, "CURRENT" if num <= 2 else "NEXT", 9.6, GREEN if num <= 2 else AMBER, True)
    rect(slide, 0.66, 5.96, 12.00, 0.72, GREEN, GREEN)
    add_text(slide, 0.90, 6.15, 11.55, 0.30, "첫 적용: 국가 안전 · 사업 정보 · 성과 설명 등 오류 비용이 높은 문장 검토", 15.2, WHITE, True, align=PP_ALIGN.CENTER)
    footer(slide)
    return slide


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = GREEN
    add_text(slide, 0.72, 0.72, 4.2, 0.30, "CLAIMGATE ODA", 11, RGBColor(192, 220, 204), True)
    add_text(slide, 0.72, 1.48, 11.88, 1.00, "AI가 초안을 빠르게 쓰게 한다면,", 28, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.72, 2.42, 11.88, 1.16, "ClaimGate ODA는 그 문장을 사람이\n근거와 함께 책임질 수 있게 합니다.", 27, WHITE, True, align=PP_ALIGN.CENTER)
    rect(slide, 2.10, 4.22, 9.15, 1.15, RGBColor(31, 96, 69), RGBColor(31, 96, 69))
    add_text(slide, 2.36, 4.49, 8.62, 0.55, "공공데이터 AI 시대의 마지막 문은 Claim Gate입니다.", 19, WHITE, True, align=PP_ALIGN.CENTER)
    pill(slide, 4.87, 5.86, 1.72, "NO ANCHOR", RGBColor(36, 97, 72), WHITE, 10)
    pill(slide, 6.72, 5.86, 1.54, "NO CLAIM", AMBER, DARK, 10)
    add_text(slide, 0.72, 6.74, 11.88, 0.32, "감사합니다", 13, RGBColor(201, 224, 211), True, align=PP_ALIGN.CENTER)
    return slide


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    for required in (ASSETS / "01-guided-demo-start.png", ASSETS / "02-review-workspace.png"):
        if not required.exists():
            raise SystemExit(f"missing presentation asset: {required}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for builder in (slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8, slide_9, slide_10):
        builder(prs)
    prs.core_properties.title = "ClaimGate ODA Final Presentation"
    prs.core_properties.subject = "2026 MOFA Public Data and AI Contest"
    prs.core_properties.author = "ClaimGate ODA"
    prs.core_properties.keywords = "ClaimGate, ODA, public data, AI, evidence"
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
